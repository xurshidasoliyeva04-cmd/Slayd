import os
import json
import sqlite3
import logging
import asyncio
import threading
from datetime import datetime
from io import BytesIO

from flask import Flask, request
from telegram import Update, InlineKeyboardButton, InlineKeyboardMarkup, Bot
from telegram.ext import (
    Application,
    CommandHandler,
    CallbackQueryHandler,
    MessageHandler,
    ContextTypes,
    filters,
)

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# 1. Logging
logging.basicConfig(format="%(asctime)s - %(name)s - %(levelname)s - %(message)s", level=logging.INFO)
logger = logging.getLogger(__name__)

# 2. Token va Sozlamalar
BOT_TOKEN = "8461901986:AAHIQLMa1RckCqGCU71PJuJZCCnfKdWjYXk"
app = Flask(__name__)

WAITING_NAME, WAITING_PAGES, WAITING_DESIGN, WAITING_TOPIC = range(4)

DESIGNS = {
    1: {"name": "🔵 Ko'k Professional", "primary": "1E3A8A", "secondary": "3B82F6", "accent": "DBEAFE", "text": "1E293B"},
    2: {"name": "🟢 Yashil Tabiat", "primary": "166534", "secondary": "22C55E", "accent": "DCFCE7", "text": "1E293B"},
    3: {"name": "🔴 Qizil Energiya", "primary": "991B1B", "secondary": "EF4444", "accent": "FEE2E2", "text": "1E293B"},
    4: {"name": "🟣 Binafsha Kreativ", "primary": "581C87", "secondary": "A855F7", "accent": "F3E8FF", "text": "1E293B"},
    5: {"name": "🟡 Sariq Quyosh", "primary": "854D0E", "secondary": "EAB308", "accent": "FEF9C3", "text": "1E293B"},
    6: {"name": "⚫ Qora Elegant", "primary": "18181B", "secondary": "3F3F46", "accent": "F4F4F5", "text": "18181B"},
}

# 3. Database Funksiyalari
def init_db():
    conn = sqlite3.connect("users.db", check_same_thread=False)
    c = conn.cursor()
    c.execute("CREATE TABLE IF NOT EXISTS users (user_id INTEGER PRIMARY KEY, username TEXT, full_name TEXT, credits INTEGER DEFAULT 2, referral_code TEXT UNIQUE, referred_by INTEGER, referral_count INTEGER DEFAULT 0, created_at TEXT)")
    c.execute("CREATE TABLE IF NOT EXISTS user_state (user_id INTEGER PRIMARY KEY, state INTEGER DEFAULT 0, data TEXT DEFAULT '{}')")
    conn.commit()
    conn.close()

def get_user(user_id):
    conn = sqlite3.connect("users.db", check_same_thread=False)
    c = conn.cursor()
    c.execute("SELECT * FROM users WHERE user_id = ?", (user_id,))
    u = c.fetchone()
    conn.close()
    return u

def create_user(user_id, username, full_name, referred_by=None):
    conn = sqlite3.connect("users.db", check_same_thread=False)
    c = conn.cursor()
    ref_code = f"REF{user_id}"
    try:
        c.execute("INSERT INTO users (user_id, username, full_name, credits, referral_code, referred_by, referral_count, created_at) VALUES (?, ?, ?, 2, ?, ?, 0, ?)", (user_id, username, full_name, ref_code, referred_by, datetime.now().isoformat()))
        if referred_by:
            c.execute("UPDATE users SET referral_count = referral_count + 1, credits = credits + 1 WHERE user_id = ?", (referred_by,))
        conn.commit()
    except: pass
    conn.close()

def get_credits(user_id):
    u = get_user(user_id)
    return u[3] if u else 0

def use_credit(user_id):
    conn = sqlite3.connect("users.db", check_same_thread=False)
    c = conn.cursor()
    c.execute("UPDATE users SET credits = credits - 1 WHERE user_id = ? AND credits > 0", (user_id,))
    success = c.rowcount > 0
    conn.commit()
    conn.close()
    return success

def get_user_state(user_id):
    conn = sqlite3.connect("users.db", check_same_thread=False)
    c = conn.cursor()
    c.execute("SELECT state, data FROM user_state WHERE user_id = ?", (user_id,))
    r = c.fetchone()
    conn.close()
    return (r[0], json.loads(r[1])) if r else (0, {})

def set_user_state(user_id, state, data=None):
    conn = sqlite3.connect("users.db", check_same_thread=False)
    c = conn.cursor()
    c.execute("INSERT OR REPLACE INTO user_state (user_id, state, data) VALUES (?, ?, ?)", (user_id, state, json.dumps(data or {})))
    conn.commit()
    conn.close()

# 4. Slayd yaratish
def hex_to_rgb(h): return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def create_presentation(topic, full_name, num_pages, design_num):
    design = DESIGNS[design_num]
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(10), Inches(5.625)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = hex_to_rgb(design["primary"])
    t_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(1))
    tf = t_box.text_frame; tf.text = topic; tf.paragraphs[0].font.size = Pt(44); tf.paragraphs[0].font.color.rgb = RGBColor(255,255,255); tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    for i in range(num_pages - 1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid(); bg.fill.fore_color.rgb = hex_to_rgb(design["accent"])
        c_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(8.5), Inches(3.5))
        c_box.text_frame.text = f"• {topic} bo'yicha ma'lumotlar {i+1}"
    buf = BytesIO(); prs.save(buf); buf.seek(0); return buf

# 5. Bot Handlers
async def start(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user = update.effective_user
    args = context.args
    ref_by = int(args[0].replace("REF", "")) if args and args[0].startswith("REF") else None
    if not get_user(user.id): create_user(user.id, user.username, user.full_name, ref_by)
    set_user_state(user.id, 0, {})
    kb = [[InlineKeyboardButton("📊 Slayd Yaratish", callback_data="create_slide")]]
    await update.message.reply_text(f"Xush kelibsiz! Kreditlar: {get_credits(user.id)}", reply_markup=InlineKeyboardMarkup(kb))

async def handle_callback(update: Update, context: ContextTypes.DEFAULT_TYPE):
    query = update.callback_query; await query.answer(); user_id = query.from_user.id
    if query.data == "create_slide":
        if get_credits(user_id) <= 0: await query.edit_message_text("Kredit qolmadi!")
        else:
            set_user_state(user_id, WAITING_NAME + 1, {})
            await query.edit_message_text("Ismingizni yozing:")
    elif query.data.startswith("pages_"):
        p = int(query.data.split("_")[1]); s, d = get_user_state(user_id); d["pages"] = p
        set_user_state(user_id, WAITING_DESIGN + 1, d)
        kb = [[InlineKeyboardButton(v["name"], callback_data=f"design_{k}")] for k, v in DESIGNS.items()]
        await query.edit_message_text("Dizayn tanlang:", reply_markup=InlineKeyboardMarkup(kb))
    elif query.data.startswith("design_"):
        ds = int(query.data.split("_")[1]); s, d = get_user_state(user_id); d["design"] = ds
        set_user_state(user_id, WAITING_TOPIC + 1, d)
        await query.edit_message_text("Mavzuni yozing:")

async def handle_message(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id; text = update.message.text; state, data = get_user_state(user_id)
    if state == WAITING_NAME + 1:
        data["full_name"] = text; set_user_state(user_id, WAITING_PAGES + 1, data)
        kb = [[InlineKeyboardButton(str(i), callback_data=f"pages_{i}") for i in range(8, 12)], [InlineKeyboardButton(str(i), callback_data=f"pages_{i}") for i in range(12, 16)]]
        await update.message.reply_text("Betlar soni:", reply_markup=InlineKeyboardMarkup(kb))
    elif state == WAITING_TOPIC + 1:
        if not use_credit(user_id): return
        await update.message.reply_text("Tayyorlanmoqda...")
        try:
            buf = create_presentation(text, data["full_name"], data["pages"], data["design"])
            await update.message.reply_document(document=buf, filename="slayd.pptx", caption="Tayyor!")
        except: await update.message.reply_text("Xato!")
        set_user_state(user_id, 0, {})

# 6. Render uchun Flask
@app.route('/')
def home(): return "Bot is live!"

def run_flask():
    port = int(os.environ.get("PORT", 8080))
    app.run(host='0.0.0.0', port=port)

# 7. Main
if __name__ == '__main__':
    init_db()
    threading.Thread(target=run_flask, daemon=True).start()
    application = Application.builder().token(BOT_TOKEN).build()
    application.add_handler(CommandHandler("start", start))
    application.add_handler(CallbackQueryHandler(handle_callback))
    application.add_handler(MessageHandler(filters.TEXT & ~filters.COMMAND, handle_message))
    application.run_polling()
